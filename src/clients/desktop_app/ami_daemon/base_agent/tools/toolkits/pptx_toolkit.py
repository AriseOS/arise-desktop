"""
PPTXToolkit - PowerPoint presentation creation for document operations.

Based on Eigent's PPTXToolkit implementation which wraps CAMEL's PPTXToolkit.
Provides PowerPoint presentation creation capabilities with event emission support.

References:
- Eigent: third-party/eigent/backend/app/utils/toolkit/pptx_toolkit.py
- CAMEL: camel.toolkits.PPTXToolkit
"""

import json
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from .base_toolkit import BaseToolkit, FunctionTool
from ...events import listen_toolkit
from ...workspace import get_working_directory

logger = logging.getLogger(__name__)


class PPTXToolkit(BaseToolkit):
    """A toolkit for creating PowerPoint presentations.

    Provides presentation creation capabilities:
    - Create presentations with title slides and content slides
    - Format text with bold and italic styling
    - Create bullet point lists
    - Create tables with headers and data
    - Support for step-by-step process slides

    Based on Eigent's implementation which wraps CAMEL's PPTXToolkit.
    """

    agent_name: str = "document_agent"

    def __init__(
        self,
        working_directory: Optional[str] = None,
        timeout: Optional[float] = 60.0,
    ) -> None:
        """Initialize the PPTXToolkit.

        Args:
            working_directory: Directory for saving presentations.
                If not provided, uses task workspace from WorkingDirectoryManager.
            timeout: Operation timeout in seconds.
        """
        super().__init__(timeout=timeout)

        # Determine working directory - fail if not provided and no workspace manager
        if working_directory:
            self._working_directory = Path(working_directory)
        else:
            self._working_directory = Path(get_working_directory())

        # Ensure directory exists
        self._working_directory.mkdir(parents=True, exist_ok=True)

        logger.info(f"PPTXToolkit initialized in {self._working_directory}")

    def _resolve_filepath(self, filename: str) -> Path:
        """Resolve a filename to an absolute path.

        Args:
            filename: The filename or path.

        Returns:
            Absolute Path object.
        """
        path = Path(filename)
        if path.is_absolute():
            return path
        return self._working_directory / filename

    @listen_toolkit(
        inputs=lambda self, content, filename, **kw: f"Creating presentation: {filename}",
        return_msg=lambda r: r[:200] if isinstance(r, str) else str(r)
    )
    def create_presentation(
        self,
        content: str,
        filename: str,
        template: Optional[str] = None,
    ) -> str:
        """Create a PowerPoint presentation.

        The content should be a JSON string representing an array of slide objects.
        Each slide can have:
        - title/subtitle: For title slides
        - heading: Slide title
        - bullet_points: List of bullet point strings
        - table: Object with 'headers' and 'rows'
        - steps: List of step descriptions for process slides

        Example content JSON:
        [
            {"title": "Main Title", "subtitle": "Subtitle"},
            {"heading": "Overview", "bullet_points": ["Point 1", "Point 2"]},
            {"heading": "Data", "table": {"headers": ["Col1", "Col2"], "rows": [["A", "B"]]}}
        ]

        Args:
            content: JSON string of slide content.
            filename: Output filename (will add .pptx if not present).
            template: Optional path to a template .pptx file.

        Returns:
            Success message with file path, or error message.
        """
        try:
            from pptx import Presentation
        except ImportError:
            return "Error: python-pptx package required. Install with: pip install python-pptx"

        # Ensure filename has .pptx extension
        if not filename.lower().endswith(".pptx"):
            filename += ".pptx"

        filepath = self._resolve_filepath(filename)
        filepath.parent.mkdir(parents=True, exist_ok=True)

        logger.info(f"Creating presentation: {filepath}")

        try:
            # Parse content JSON
            if isinstance(content, str):
                slides_data = json.loads(content)
            else:
                slides_data = content

            if not isinstance(slides_data, list):
                return "Error: Content must be a JSON array of slide objects"

            # Create presentation from template or blank
            if template and Path(template).exists():
                prs = Presentation(template)
            else:
                prs = Presentation()

            # Process each slide
            for slide_data in slides_data:
                self._add_slide(prs, slide_data)

            # Save presentation
            prs.save(str(filepath))

            logger.info(f"PowerPoint presentation successfully created: {filepath}")
            return f"PowerPoint presentation successfully created: {filepath}"

        except json.JSONDecodeError as e:
            error_msg = f"Error parsing content JSON: {str(e)}"
            logger.error(error_msg)
            return error_msg
        except Exception as e:
            error_msg = f"Error creating presentation: {str(e)}"
            logger.error(error_msg)
            return error_msg

    def _add_slide(self, prs: Any, slide_data: Dict[str, Any]) -> None:
        """Add a slide to the presentation.

        Args:
            prs: Presentation object.
            slide_data: Dictionary with slide content.
        """
        from pptx.util import Inches, Pt

        # Determine slide type based on content
        if "title" in slide_data and ("subtitle" in slide_data or len(slide_data) <= 2):
            # Title slide
            layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", "")

            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = slide_data.get("subtitle", "")

        elif "table" in slide_data:
            # Table slide
            layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(layout)

            # Add heading
            if "heading" in slide_data:
                self._add_title_shape(slide, slide_data["heading"])

            # Add table
            table_data = slide_data["table"]
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])

            if headers or rows:
                num_cols = len(headers) if headers else len(rows[0]) if rows else 1
                num_rows = len(rows) + (1 if headers else 0)

                # Create table
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(0.5 * num_rows)

                table = slide.shapes.add_table(
                    num_rows, num_cols, left, top, width, height
                ).table

                # Add headers
                if headers:
                    for col_idx, header in enumerate(headers):
                        cell = table.cell(0, col_idx)
                        cell.text = str(header)

                # Add data rows
                start_row = 1 if headers else 0
                for row_idx, row in enumerate(rows):
                    for col_idx, cell_value in enumerate(row):
                        if col_idx < num_cols:
                            cell = table.cell(start_row + row_idx, col_idx)
                            cell.text = str(cell_value)

        elif "bullet_points" in slide_data:
            # Bullet points slide
            layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("heading", "")

            # Find content placeholder
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break

            if body_shape:
                tf = body_shape.text_frame
                tf.clear()

                bullet_points = slide_data.get("bullet_points", [])
                for i, point in enumerate(bullet_points):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(point)
                    p.level = 0

        elif "steps" in slide_data:
            # Process/steps slide
            layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(layout)

            if "heading" in slide_data:
                self._add_title_shape(slide, slide_data["heading"])

            steps = slide_data.get("steps", [])
            for i, step in enumerate(steps):
                left = Inches(0.5 + (i % 3) * 3)
                top = Inches(2 + (i // 3) * 1.5)
                width = Inches(2.5)
                height = Inches(1)

                shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    left, top, width, height
                )
                shape.text = f"{i + 1}. {step}"

        else:
            # Default content slide
            layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title and "heading" in slide_data:
                slide.shapes.title.text = slide_data["heading"]

            # Add any text content
            if "content" in slide_data:
                body_shape = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        body_shape = shape
                        break

                if body_shape:
                    body_shape.text = slide_data["content"]

    def _add_title_shape(self, slide: Any, title: str) -> None:
        """Add a title text box to a slide.

        Args:
            slide: Slide object.
            title: Title text.
        """
        from pptx.util import Inches, Pt

        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.75)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True

    @property
    def working_directory(self) -> Path:
        """Get the current working directory."""
        return self._working_directory

    def get_tools(self) -> List[FunctionTool]:
        """Return a list of FunctionTool objects for this toolkit.

        Returns:
            List of FunctionTool objects.
        """
        return [
            FunctionTool(self.create_presentation),
        ]

    @classmethod
    def toolkit_name(cls) -> str:
        """Return the name of this toolkit."""
        return "PPTX"
